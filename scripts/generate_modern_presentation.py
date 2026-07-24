import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Palette de couleurs ---
COLORS = {
    "Fond principal": RGBColor(247, 244, 238),       # #F7F4EE
    "Surface": RGBColor(255, 255, 255),              # #FFFFFF
    "Texte principal": RGBColor(37, 34, 31),         # #25221F
    "Texte secondaire": RGBColor(107, 101, 95),      # #6B655F
    "Terracotta principal": RGBColor(201, 100, 66),  # #C96442
    "Terracotta clair": RGBColor(241, 216, 206),     # #F1D8CE
    "Vert secondaire": RGBColor(61, 122, 104),       # #3D7A68
    "Vert clair": RGBColor(220, 235, 229),           # #DCEBE5
    "Bleu technique": RGBColor(73, 106, 129),        # #496A81
    "Bordure": RGBColor(222, 216, 207),              # #DED8CF
    "Erreur": RGBColor(184, 74, 74)                  # #B84A4A
}

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["Fond principal"]

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["Texte principal"]
    p.font.name = "Arial"
    return title_box

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def draw_card(slide, x, y, w, h, title, text, bg_color="Surface", text_color="Texte secondaire"):
    # Ombre légère n'est pas nativement supportée par python-pptx sans manipulation XML
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[bg_color]
    shape.line.color.rgb = COLORS["Bordure"]
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    
    if title:
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(20)
        p1.font.color.rgb = COLORS["Texte principal"]
        p1.alignment = PP_ALIGN.CENTER
        
    if text:
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLORS[text_color]
        p2.alignment = PP_ALIGN.CENTER
    return shape

def draw_arrow(slide, x1, y1, x2, y2, color="Texte secondaire"):
    # Approximer une flèche avec un connecteur droit
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = COLORS[color]
    shape.line.width = Pt(3)
    # Ajouter la pointe de flèche via XML n'est pas trivial, on utilise une ligne simple
    return shape

def main():
    prs = Presentation()
    # Format 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank layout
    
    # ---------------------------------------------------------
    # SLIDE 1 : Couverture
    # ---------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)
    
    # Textes de gauche
    tb = s1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(5), Inches(3))
    tf = tb.text_frame
    p = tf.add_paragraph()
    p.text = "YouCode AI Platform"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = COLORS["Terracotta principal"]
    
    p2 = tf.add_paragraph()
    p2.text = "Plateforme conversationnelle multi-agent"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLORS["Texte secondaire"]
    
    p3 = tf.add_paragraph()
    p3.text = "\nAziz BENMALLOUK"
    p3.font.size = Pt(20)
    p3.font.color.rgb = COLORS["Texte principal"]
    
    # Visuel abstrait de droite (Noyau + WhatsApp + Agents)
    draw_card(s1, 8, 3, 2, 1.5, "Orchestrator", "", "Terracotta clair")
    draw_card(s1, 11, 1.5, 1.5, 1, "Guide", "", "Vert clair")
    draw_card(s1, 11, 3.25, 1.5, 1, "Support", "", "Vert clair")
    draw_card(s1, 11, 5, 1.5, 1, "Newsletter", "", "Vert clair")
    draw_card(s1, 5.5, 3.25, 1.5, 1, "WhatsApp", "", "Surface")
    
    draw_arrow(s1, 7, 3.75, 8, 3.75)
    draw_arrow(s1, 10, 3.75, 11, 2)
    draw_arrow(s1, 10, 3.75, 11, 3.75)
    draw_arrow(s1, 10, 3.75, 11, 5.5)
    
    add_speaker_notes(s1, "Bonjour à tous. Je vous présente YouCode AI Platform, un système d'assistance automatisée. Le but est de réduire la charge manuelle tout en gardant une interface simple.")

    # ---------------------------------------------------------
    # SLIDE 2 : Problématique
    # ---------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_title(s2, "Le Besoin")
    
    # Gauche : Problèmes
    draw_card(s2, 1, 2, 4, 0.8, "", "Informations dispersées")
    draw_card(s2, 1, 3.2, 4, 0.8, "", "Demandes répétitives")
    draw_card(s2, 1, 4.4, 4, 0.8, "", "Suivi manuel")
    draw_card(s2, 1, 5.6, 4, 0.8, "", "Disponibilité limitée")
    
    # Flèche
    draw_arrow(s2, 5.5, 4, 7, 4, "Terracotta principal")
    
    # Droite : Solution
    draw_card(s2, 7.5, 3, 5, 2, "Assistance fiable", "Multilingue et disponible sur WhatsApp", "Terracotta clair", "Texte principal")
    
    add_speaker_notes(s2, "Nous faisons face à de multiples requêtes répétitives. La solution doit être accessible là où sont nos candidats : sur WhatsApp.")

    # ---------------------------------------------------------
    # SLIDE 3 : Solution
    # ---------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_title(s3, "Flux de Résolution")
    
    draw_card(s3, 0.5, 3, 2, 1.5, "WhatsApp", "Canal utilisateur")
    draw_arrow(s3, 2.5, 3.75, 3, 3.75)
    
    draw_card(s3, 3, 3, 2, 1.5, "FastAPI", "Point d'entrée")
    draw_arrow(s3, 5, 3.75, 5.5, 3.75)
    
    draw_card(s3, 5.5, 3, 2, 1.5, "Multi-Agent", "Orchestration LangGraph")
    draw_arrow(s3, 7.5, 3.75, 8, 3.75)
    
    draw_card(s3, 8, 3, 2.5, 1.5, "RAG & MCP", "Services externes (Qdrant, Sheets)")
    draw_arrow(s3, 10.5, 3.75, 11, 3.75)
    
    draw_card(s3, 11, 3, 2, 1.5, "Réponse", "Personnalisée")
    
    add_speaker_notes(s3, "L'architecture s'appuie sur une pipeline allant de WhatsApp à FastAPI, qui distribue ensuite la tâche aux agents via LangGraph. Ces agents consomment le RAG et écrivent dans Sheets.")

    # ---------------------------------------------------------
    # SLIDE 4 : Cas d'utilisation
    # ---------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_title(s4, "Cas d'Utilisation Couverts")
    
    # Centre
    draw_card(s4, 5.5, 3, 2.33, 1.5, "Visiteur", "Centre d'attention", "Vert clair")
    
    # Branches
    draw_card(s4, 2, 1.5, 2.5, 1, "", "Découvrir YouCode")
    draw_card(s4, 9, 1.5, 2.5, 1, "", "Report de test")
    draw_card(s4, 2, 5, 2.5, 1, "", "Support")
    draw_card(s4, 9, 5, 2.5, 1, "", "Newsletter")
    
    # Connecteurs
    draw_arrow(s4, 4.5, 2, 5.5, 3.5)
    draw_arrow(s4, 9, 2, 7.83, 3.5)
    draw_arrow(s4, 4.5, 5.5, 5.5, 4)
    draw_arrow(s4, 9, 5.5, 7.83, 4)
    
    add_speaker_notes(s4, "Le système gère plusieurs intentions : découverte, démarches administratives, abonnement, report d'admission.")

    # ---------------------------------------------------------
    # SLIDE 5 : Architecture Globale
    # ---------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_title(s5, "Architecture en Couches")
    
    draw_card(s5, 2, 1.5, 9, 0.8, "Canal", "WhatsApp", "Surface")
    draw_card(s5, 2, 2.6, 9, 0.8, "Entrée", "Webhook FastAPI", "Surface")
    draw_card(s5, 2, 3.7, 9, 0.8, "Orchestration", "LangGraph + Supervisor", "Terracotta clair")
    
    # Agents
    draw_card(s5, 2, 4.8, 2.8, 0.8, "Guide", "Agent", "Vert clair")
    draw_card(s5, 5.1, 4.8, 2.8, 0.8, "Support", "Agent", "Vert clair")
    draw_card(s5, 8.2, 4.8, 2.8, 0.8, "Newsletter", "Agent", "Vert clair")
    
    draw_card(s5, 2, 5.9, 9, 0.8, "Données & Services", "Qdrant | MCP | Google Sheets", "Bleu technique", "Surface")
    
    add_speaker_notes(s5, "Une architecture modulaire : le webhook FastAPI passe le relais au Supervisor LangGraph, qui délègue ensuite aux agents. Chaque agent a accès à ses propres services de données.")

    # ---------------------------------------------------------
    # SLIDE 6 : Architecture Multi-Agent
    # ---------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_title(s6, "Spécialisation des Agents")
    
    draw_card(s6, 5.5, 3, 2.33, 1.5, "Supervisor", "Routage", "Terracotta clair")
    
    draw_card(s6, 1, 1.5, 3.5, 1.5, "Guide Agent", "Rôle : Informations\nOutil : RAG (Qdrant)", "Vert clair")
    draw_card(s6, 9, 1.5, 3.5, 1.5, "Support Agent", "Rôle : Démarches\nOutil : Formulaires", "Vert clair")
    draw_card(s6, 5, 5.5, 3.5, 1.5, "Newsletter Agent", "Rôle : Inscription\nOutil : Sheets via MCP", "Vert clair")
    
    draw_arrow(s6, 4.5, 2.25, 5.5, 3.5)
    draw_arrow(s6, 9, 2.25, 7.83, 3.5)
    draw_arrow(s6, 6.66, 5.5, 6.66, 4.5)
    
    add_speaker_notes(s6, "Chaque agent a une responsabilité unique. Le Guide utilise le RAG pour ne pas halluciner. Le Support collecte des entités pour les démarches.")

    # ---------------------------------------------------------
    # SLIDE 7 : Routage et Continuité
    # ---------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_title(s7, "Maintien du Contexte")
    
    # Gauche
    draw_card(s7, 1, 2, 4, 3, "Nouvelle intention", "START\n↓\nSupervisor\n↓\nNouvel Agent", "Surface")
    
    # Droite
    draw_card(s7, 8, 2, 4, 3, "Conversation active", "START\n↓\nPhase active\n↓\nAgent actuel", "Terracotta clair")
    
    # Bas
    draw_card(s7, 1, 5.5, 11, 1, "", "Un workflow actif conserve la conversation sans nouvelle classification.", "Vert clair")
    
    add_speaker_notes(s7, "LangGraph garde l'état de la conversation. Si le candidat remplit déjà un formulaire, on by-pass la détection d'intention pour retourner directement dans la phase active.")

    # ---------------------------------------------------------
    # SLIDE 8 : Graph LangGraph
    # ---------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_title(s8, "Graphe d'État LangGraph")
    
    draw_card(s8, 5, 1.5, 3.33, 0.8, "Supervisor", "Classification")
    
    # Support
    draw_card(s8, 1, 3, 3, 0.6, "", "Support: Extract")
    draw_card(s8, 1, 4, 3, 0.6, "", "Support: Consent")
    draw_card(s8, 1, 5, 3, 0.6, "", "Support: Process")
    draw_card(s8, 1, 6, 3, 0.6, "", "Support: Decision")
    
    # Newsletter
    draw_card(s8, 9, 3, 3, 0.6, "", "News: Extract")
    draw_card(s8, 9, 4, 3, 0.6, "", "News: Consent")
    draw_card(s8, 9, 5, 3, 0.6, "", "News: Process")
    
    # END
    draw_card(s8, 5, 6, 3.33, 0.8, "END", "Arrêt / Reprise", "Terracotta clair")
    
    # Arrow hints
    draw_arrow(s8, 5, 1.9, 2.5, 3)
    draw_arrow(s8, 8.33, 1.9, 10.5, 3)
    
    add_speaker_notes(s8, "Ceci est le graphe réel. Les noeuds s'exécutent de manière séquentielle, mais peuvent atteindre la fin (END) prématurément pour attendre l'input WhatsApp suivant.")

    # ---------------------------------------------------------
    # SLIDE 9 : Pipeline RAG
    # ---------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_title(s9, "RAG : Parent-Child Strategy")
    
    # Pipeline
    y_pos = 3
    draw_card(s9, 0.5, y_pos, 1.5, 1, "", "Question")
    draw_arrow(s9, 2, y_pos+0.5, 2.5, y_pos+0.5)
    draw_card(s9, 2.5, y_pos, 2, 1, "", "Recherche Child\n(Petit chunk)")
    draw_arrow(s9, 4.5, y_pos+0.5, 5, y_pos+0.5)
    draw_card(s9, 5, y_pos, 2, 1, "", "Parent Context\n(Doc complet)", "Terracotta clair")
    draw_arrow(s9, 7, y_pos+0.5, 7.5, y_pos+0.5)
    draw_card(s9, 7.5, y_pos, 2, 1, "", "LLM\n(Génération)")
    draw_arrow(s9, 9.5, y_pos+0.5, 10, y_pos+0.5)
    draw_card(s9, 10, y_pos, 2.8, 1, "", "Réponse vérifiée")
    
    # Garanties
    draw_card(s9, 2, 5.5, 2.5, 0.8, "", "Docs officiels", "Vert clair")
    draw_card(s9, 5.5, 5.5, 2.5, 0.8, "", "Multilingue", "Vert clair")
    draw_card(s9, 9, 5.5, 2.5, 0.8, "", "Zéro invention", "Vert clair")
    
    add_speaker_notes(s9, "Le retriever utilise Qdrant. La stratégie Parent-Child permet de trouver l'info exacte dans un petit bout de texte, puis de charger le document complet pour le LLM.")

    # ---------------------------------------------------------
    # SLIDE 10 : Support Conversationnel
    # ---------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    apply_background(s10)
    add_title(s10, "Timeline du Support")
    
    # Timeline
    draw_card(s10, 1, 2, 5, 0.6, "", "1. Détection demande")
    draw_card(s10, 1, 3, 5, 0.6, "", "2. Extraction structurée", "Terracotta clair")
    draw_card(s10, 1, 4, 5, 0.6, "", "3. Question au visiteur (champ manquant)")
    draw_card(s10, 1, 5, 5, 0.6, "", "4. Validation du consentement")
    draw_card(s10, 1, 6, 5, 0.6, "", "5. Enregistrement")
    
    # State Card
    draw_card(s10, 8, 3, 4, 2, "State Interne", "active_agent : support\nphase : extraction\ndraft : {email: ...}", "Surface")
    
    add_speaker_notes(s10, "L'agent Support utilise Pydantic pour l'extraction structurée. S'il manque un email, l'agent pose automatiquement la question et met le graphe en pause.")

    # ---------------------------------------------------------
    # SLIDE 11 : Diagramme de Séquence (Report de test)
    # ---------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    apply_background(s11)
    add_title(s11, "Flux métier : Report de Test")
    
    # Lignes de vie
    draw_card(s11, 1, 1.5, 2, 0.8, "", "Visiteur")
    draw_arrow(s11, 2, 2.3, 2, 6.5)
    
    draw_card(s11, 4, 1.5, 2, 0.8, "", "Support Agent")
    draw_arrow(s11, 5, 2.3, 5, 6.5)
    
    draw_card(s11, 7, 1.5, 2, 0.8, "", "MCP Sheets", "Terracotta clair")
    draw_arrow(s11, 8, 2.3, 8, 6.5)
    
    draw_card(s11, 10, 1.5, 2, 0.8, "", "Responsable")
    draw_arrow(s11, 11, 2.3, 11, 6.5)
    
    # Actions (simples textes)
    draw_card(s11, 2.5, 3, 2, 0.4, "", "1. Demande")
    draw_arrow(s11, 2, 3.2, 5, 3.2)
    
    draw_card(s11, 5.5, 4, 2, 0.4, "", "2. Brouillon")
    draw_arrow(s11, 5, 4.2, 8, 4.2)
    
    draw_card(s11, 2.5, 5, 2, 0.4, "", "3. Proposition")
    draw_arrow(s11, 5, 5.2, 2, 5.2)
    
    draw_card(s11, 8.5, 6, 2, 0.4, "", "4. Validation")
    draw_arrow(s11, 8, 6.2, 11, 6.2)
    
    add_speaker_notes(s11, "Pour un report de test : le visiteur demande, l'agent crée un brouillon via MCP, propose une date, et finalement la demande arrive au statut 'À Valider' pour le Responsable.")

    # ---------------------------------------------------------
    # SLIDE 12 : Choix d'une session
    # ---------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    apply_background(s12)
    add_title(s12, "Prise de décision interactive")
    
    # Decision tree
    draw_card(s12, 1, 2, 4, 1, "", "Session proposée")
    draw_arrow(s12, 3, 3, 3, 4)
    draw_card(s12, 1, 4, 4, 1, "", "Convient-elle ?")
    
    draw_arrow(s12, 3, 5, 2, 6)
    draw_arrow(s12, 3, 5, 4, 6)
    
    draw_card(s12, 0.5, 6, 2, 1, "", "Oui → Valider", "Vert clair")
    draw_card(s12, 3.5, 6, 2, 1, "", "Non → Alternative", "Terracotta clair")
    
    # Mock WhatsApp
    draw_card(s12, 8, 2, 4, 4.5, "WhatsApp", "Agent : Session 20 août ?\n\nVisiteur : Non, une autre.", "Surface")
    
    add_speaker_notes(s12, "L'agent ne force pas les choix. Il interroge les APIs pour trouver des créneaux, les propose au candidat et adapte son comportement selon la réponse.")

    # ---------------------------------------------------------
    # SLIDE 13 : MCP et Google Sheets
    # ---------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    apply_background(s13)
    add_title(s13, "Model Context Protocol (MCP)")
    
    draw_card(s13, 0.5, 2, 2.5, 1.2, "", "Agent")
    draw_arrow(s13, 3, 2.6, 3.5, 2.6)
    draw_card(s13, 3.5, 2, 2.5, 1.2, "", "Client MCP\n(Standard)")
    draw_arrow(s13, 6, 2.6, 6.5, 2.6)
    draw_card(s13, 6.5, 2, 2.5, 1.2, "", "Serveur MCP\n(FastMCP)")
    draw_arrow(s13, 9, 2.6, 9.5, 2.6)
    draw_card(s13, 9.5, 2, 3, 1.2, "", "Google Sheets", "Vert clair")
    
    # Tableau comparatif
    draw_card(s13, 1, 4.5, 5, 2, "State (SQLite)", "Mémoire Temporaire\nBrouillons", "Surface")
    draw_card(s13, 7, 4.5, 5, 2, "Google Sheets", "Mémoire Persistante\nSuivi Métier", "Terracotta clair")
    
    add_speaker_notes(s13, "MCP standardise la communication avec les outils. Les agents ne savent pas qu'ils utilisent Sheets, ils appellent juste des méthodes unifiées. Cela permet un découplage total.")

    # ---------------------------------------------------------
    # SLIDE 14 : FastAPI Intégrations
    # ---------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    apply_background(s14)
    add_title(s14, "Moteur Central : FastAPI")
    
    draw_card(s14, 5.5, 3, 2.5, 1.5, "FastAPI", "Webhook & API", "Bleu technique", "Surface")
    
    draw_card(s14, 2, 1.5, 2.5, 1, "", "WhatsApp Webhook")
    draw_card(s14, 9, 1.5, 2.5, 1, "", "Qdrant Vector DB")
    draw_card(s14, 2, 5, 2.5, 1, "", "LangGraph Graph")
    draw_card(s14, 9, 5, 2.5, 1, "", "Google Sheets MCP")
    
    draw_arrow(s14, 4.5, 2.5, 5.5, 3.75)
    draw_arrow(s14, 9, 2.5, 8, 3.75)
    draw_arrow(s14, 4.5, 5.5, 5.5, 3.75)
    draw_arrow(s14, 9, 5.5, 8, 3.75)
    
    add_speaker_notes(s14, "FastAPI est la clé de voûte. Il gère l'asynchronisme pour ne pas bloquer les messages WhatsApp pendant que le LLM génère sa réponse.")

    # ---------------------------------------------------------
    # SLIDE 15 : Sécurité et Fiabilité
    # ---------------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    apply_background(s15)
    add_title(s15, "Garanties Fonctionnelles")
    
    draw_card(s15, 0.5, 2.5, 2.5, 3, "Consentement", "Validation explicite requise.", "Surface")
    draw_card(s15, 3.5, 2.5, 2.5, 3, "Validation", "Typage fort Pydantic.", "Surface")
    draw_card(s15, 6.5, 2.5, 2.5, 3, "Outils", "Sécurisation MCP.", "Surface")
    draw_card(s15, 9.5, 2.5, 2.5, 3, "Humain", "Responsable dans la boucle.", "Surface")
    
    # Ligne défensive
    draw_card(s15, 0.5, 6, 11.5, 0.8, "", "Défense en profondeur", "Terracotta clair")
    
    add_speaker_notes(s15, "Pas de validation de date sans humain. Pas d'insertion sans consentement explicite. Tout est cadré pour protéger YouCode et le visiteur.")

    # ---------------------------------------------------------
    # SLIDE 16 : Chiffres et KPIs
    # ---------------------------------------------------------
    s16 = prs.slides.add_slide(blank_layout)
    apply_background(s16)
    add_title(s16, "Métrique du Projet")
    
    draw_card(s16, 1, 2, 3, 2, "3", "Agents Spécialisés", "Vert clair")
    draw_card(s16, 5, 2, 3, 2, "2", "Workflows Métier", "Terracotta clair")
    draw_card(s16, 9, 2, 3, 2, "1", "Serveur MCP", "Bleu technique", "Surface")
    
    draw_card(s16, 1, 4.5, 11, 1.5, "Bénéfices", "Disponibilité 24/7 | Automatisation fiable | Traçabilité Sheets", "Surface")
    
    add_speaker_notes(s16, "Les chiffres basés sur l'implémentation actuelle prouvent que le noyau est minimaliste mais couvre 100% de la surface attendue pour les démarches de base.")

    # ---------------------------------------------------------
    # SLIDE 17 : Roadmap
    # ---------------------------------------------------------
    s17 = prs.slides.add_slide(blank_layout)
    apply_background(s17)
    add_title(s17, "Limites et Perspectives")
    
    draw_card(s17, 1, 3, 3, 3, "Aujourd'hui", "WhatsApp\nRAG\nMCP Sheets", "Vert clair")
    draw_card(s17, 5, 3, 3, 3, "Court Terme", "Interface Admin\nAgent Guardrails", "Surface")
    draw_card(s17, 9, 3, 3, 3, "Long Terme", "PostgreSQL\nAnalytics Avancés", "Terracotta clair")
    
    draw_arrow(s17, 4, 4.5, 5, 4.5)
    draw_arrow(s17, 8, 4.5, 9, 4.5)
    
    add_speaker_notes(s17, "Nous avons livré le MVP complet. La prochaine étape est de remplacer Google Sheets par une vraie interface d'administration en Next.js.")

    # ---------------------------------------------------------
    # SLIDE 18 : Conclusion
    # ---------------------------------------------------------
    s18 = prs.slides.add_slide(blank_layout)
    apply_background(s18)
    add_title(s18, "L'essentiel")
    
    draw_card(s18, 4, 2, 5, 1, "", "Point d'entrée simple (WhatsApp)", "Surface")
    draw_card(s18, 4, 3.5, 5, 1, "", "Intelligence organisée (Agents)", "Surface")
    draw_card(s18, 4, 5, 5, 1, "", "Exécution contrôlée (MCP & Humain)", "Surface")
    
    tb = s18.shapes.add_textbox(Inches(4), Inches(6.5), Inches(5), Inches(1))
    p = tb.text_frame.add_paragraph()
    p.text = "Merci pour votre attention.\nQuestions ?"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    
    add_speaker_notes(s18, "Merci pour votre attention. Je suis maintenant à votre disposition pour la démo ou vos questions techniques.")

    # Save
    import os
    os.makedirs("docs/presentation", exist_ok=True)
    prs.save("docs/presentation/YouCode_AI_Platform_Soutenance_Modern.pptx")
    print("Présentation visuelle générée.")

if __name__ == "__main__":
    main()
