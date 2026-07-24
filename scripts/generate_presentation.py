import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Couleurs du thème YouCode (simulées)
    terracotta = RGBColor(204, 85, 0)
    anthracite = RGBColor(47, 79, 79)
    green = RGBColor(34, 139, 34)

    # Configuration des layouts
    title_slide_layout = prs.slide_layouts[0]
    title_content_layout = prs.slide_layouts[1]
    
    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = terracotta
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = anthracite

    def add_content_slide(title_text, content_items):
        slide = prs.slides.add_slide(title_content_layout)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = terracotta
        title.text_frame.paragraphs[0].font.bold = True
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        if isinstance(content_items, str):
            tf.text = content_items
        else:
            tf.text = content_items[0]
            for item in content_items[1:]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.color.rgb = anthracite

    # Slide 1 : Page de garde
    add_title_slide(
        "YouCode AI Platform",
        "Plateforme conversationnelle multi-agent\n"
        "pour l’accompagnement des visiteurs et candidats YouCode\n\n"
        "Présenté par : Aziz BENMALLOUK\n"
        "Technologies : FastAPI, LangGraph, Qdrant, MCP, WhatsApp"
    )

    # Slide 2 : Contexte
    add_content_slide(
        "Contexte",
        [
            "Sollicitations massives des visiteurs et candidats.",
            "Informations officielles dispersées et parfois obsolètes.",
            "Questions très répétitives (campus, pédagogie, admission).",
            "Démarches parfois complexes (report de test, suivi).",
            "Public multilingue (français, arabe, darija).",
            "Traitement manuel lent et fastidieux."
        ]
    )

    # Slide 3 : Problématique
    add_content_slide(
        "Problématique",
        "Comment proposer aux visiteurs de YouCode une assistance "
        "fiable, multilingue et disponible via WhatsApp, tout en "
        "automatisant les workflows répétitifs et en conservant une "
        "validation humaine pour les décisions sensibles ?"
    )

    # Slide 4 : Solution proposée
    add_content_slide(
        "Solution proposée",
        [
            "WhatsApp : canal utilisateur principal et universel.",
            "FastAPI : backend et webhook haute performance.",
            "LangGraph : orchestration dynamique et robuste (State).",
            "Agents Spécialisés : Supervisor, Guide, Support, Newsletter.",
            "RAG (Qdrant) : garantie de réponses basées sur les documents officiels.",
            "MCP (Google Sheets) : standardisation et persistance métier.",
            "Validation Humaine : pour les décisions critiques."
        ]
    )

    # Slide 5 : Utilisateurs et cas d'utilisation
    add_content_slide(
        "Utilisateurs et cas d'utilisation",
        [
            "Visiteur : Découvrir YouCode, poser des questions générales.",
            "Candidat : Signaler un problème, demander un report de test.",
            "Responsable YouCode : Consulter les demandes, valider les reports."
        ]
    )

    # Slide 6 : Architecture globale
    add_content_slide(
        "Architecture globale",
        "Visiteur (WhatsApp) -> Client WhatsApp -> Webhook FastAPI\n\n"
        "FastAPI -> LangGraph (Supervisor)\n\n"
        "Supervisor -> Agents (Guide, Support, Newsletter)\n\n"
        "Agents -> RAG (Qdrant) / APIs / MCP (Google Sheets)\n\n"
        "Retour WhatsApp"
    )

    # Slide 7 : Architecture multi-agent
    add_content_slide(
        "Architecture multi-agent",
        [
            "1. Supervisor Agent : Détecte l'intention et route la conversation.",
            "2. Guide Agent : Expert métier. Utilise le RAG.",
            "3. Support Agent : Gère les demandes (ex: reports). Collecte structurée.",
            "4. Newsletter Agent : Gère les abonnements. Demande le consentement."
        ]
    )

    # Slide 8 : Supervisor et routage
    add_content_slide(
        "Supervisor et routage",
        [
            "Nouvelle intention : Supervisor analyse -> Agent spécialiste.",
            "Workflow actif : Si l'utilisateur est dans une boucle 'Support', "
            "le message va directement au noeud actif.",
            "State géré : active_agent, support_phase, newsletter_phase, messages."
        ]
    )

    # Slide 9 : Graph LangGraph
    add_content_slide(
        "Orchestration avec LangGraph",
        [
            "Nodes (noeuds) : Exécution d'une logique ou d'un LLM.",
            "Edges (liens) : Conditions de passage au noeud suivant.",
            "State (état) : Mémoire partagée de la conversation.",
            "Capacité d'interruption : Le graphe s'arrête pour attendre "
            "une réponse WhatsApp et reprend grâce au State."
        ]
    )

    # Slide 10 : Guide Agent et RAG
    add_content_slide(
        "Guide Agent et RAG",
        [
            "Reformulation autonome de la question WhatsApp.",
            "Recherche Qdrant (Base Vectorielle).",
            "Architecture Parent-Child : Trouver le sous-chunk exact, "
            "renvoyer le contexte complet (Parent).",
            "Règle d'or : Prévention stricte des hallucinations."
        ]
    )

    # Slide 11 : Support Agent
    add_content_slide(
        "Support Agent",
        [
            "Extraction Structurée (Pydantic).",
            "Vérification des champs obligatoires (ex: campus, e-mail).",
            "Questionnement progressif de l'utilisateur.",
            "Consentement explicite exigé.",
            "Écriture via MCP dans Google Sheets."
        ]
    )

    # Slide 12 : Cas détaillé : Report de test
    add_content_slide(
        "Cas détaillé : Report de test",
        [
            "1. Collecte des données (Support Agent).",
            "2. Consentement du candidat.",
            "3. Brouillon enregistré (Google Sheets via MCP).",
            "4. Recherche d'alternative (APIs).",
            "5. Validation du candidat.",
            "6. Demande de validation finale au responsable humain."
        ]
    )

    # Slide 13 : Newsletter Agent
    add_content_slide(
        "Newsletter Agent",
        [
            "Collecte progressive (e-mail, préférences, etc.).",
            "Interruption LangGraph pour demander consentement.",
            "Analyse de la réponse au consentement.",
            "Création de l'abonnement via MCP."
        ]
    )

    # Slide 14 : MCP et Google Sheets
    add_content_slide(
        "MCP et Google Sheets",
        [
            "MCP (Model Context Protocol) : standardise les outils.",
            "Serveur MCP : Expose `append_row` ou `read_sheet`.",
            "Client MCP : Permet aux Agents d'agir sans connaître la logique métier.",
            "Avantages : Découplage, sécurité, remplacement facile du backend."
        ]
    )

    # Slide 15 : Architecture des données
    add_content_slide(
        "Architecture des données",
        [
            "State LangGraph (SQLite) : Mémoire à court terme "
            "(brouillons, suivi conversationnel).",
            "Persistance Métier (Google Sheets via MCP) : Base officielle "
            "à long terme pour les équipes YouCode.",
            "Synchronisation : Événements SQLAlchemy asynchrones."
        ]
    )

    # Slide 16 : FastAPI et intégrations
    add_content_slide(
        "FastAPI et Intégrations",
        [
            "Webhook WhatsApp (Réception asynchrone).",
            "Endpoints Métier (REST pour futures interfaces).",
            "Gestion des erreurs, Timeout et Documentation OpenAPI.",
            "Moteur de la communication avec LangGraph et MCP."
        ]
    )

    # Slide 17 : Sécurité, limites et fiabilité
    add_content_slide(
        "Sécurité, limites et fiabilité",
        [
            "Sécurité : Consentement strict, RAG basé sur documents officiels.",
            "Fiabilité : Outils fortement typés (Pydantic), validation des entités.",
            "Limites : Dépendance au réseau (WhatsApp, LLM, APIs).",
            "Parade : Timeout FastAPI et validation humaine des décisions lourdes."
        ]
    )

    # Slide 18 : Résultats et perspectives
    add_content_slide(
        "Résultats et Perspectives",
        [
            "Opérationnel : WhatsApp, RAG Qdrant, MCP Sheets, Orchestration LLM.",
            "Partiel : Modération proactive (Guardrails).",
            "Perspectives : Dashboard Administrateur web dédié, ",
            "Migration totale vers PostgreSQL, Déploiement Cloud.",
            "\nMerci pour votre attention !"
        ]
    )

    prs.save("docs/presentation/YouCode_AI_Platform_Soutenance.pptx")
    print("Présentation générée avec succès : docs/presentation/YouCode_AI_Platform_Soutenance.pptx")

if __name__ == "__main__":
    create_presentation()
