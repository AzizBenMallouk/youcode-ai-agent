from pptx import Presentation
import os

def word_count_in_shape(shape):
    if not shape.has_text_frame:
        return 0
    words = 0
    for paragraph in shape.text_frame.paragraphs:
        words += len(paragraph.text.split())
    return words

def check_presentation(pptx_path):
    prs = Presentation(pptx_path)
    
    report_lines = ["# Rapport d'Assurance Qualité Visuelle (QA)", ""]
    report_lines.append(f"**Fichier analysé :** `{os.path.basename(pptx_path)}`")
    report_lines.append("")
    
    all_passed = True
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        words_on_slide = 0
        has_large_paragraph = False
        
        for shape in slide.shapes:
            words = word_count_in_shape(shape)
            words_on_slide += words
            
            # Check for paragraphs
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if len(paragraph.text.split()) > 20: # Un peu plus de 20 mots dans un seul bloc = paragraphe suspect
                        has_large_paragraph = True
        
        status = "✅ PASS"
        issues = []
        
        if words_on_slide > 45:
            status = "❌ FAIL"
            issues.append(f"Trop de mots ({words_on_slide} > 45)")
            all_passed = False
            
        if has_large_paragraph:
            status = "❌ FAIL"
            issues.append("Présence d'un paragraphe dense détectée")
            all_passed = False
            
        report_lines.append(f"### Slide {slide_num} : {status}")
        report_lines.append(f"- **Nombre de mots visibles :** {words_on_slide}")
        if issues:
            report_lines.append(f"- **Problèmes :** {', '.join(issues)}")
        report_lines.append("")
        
    if all_passed:
        report_lines.append("## Résultat Global : SUCCÈS")
        report_lines.append("La présentation respecte les contraintes strictes du storyboard visual-first.")
    else:
        report_lines.append("## Résultat Global : ÉCHEC")
        report_lines.append("Certaines slides nécessitent une correction pour respecter la règle des 45 mots max.")
        
    # Write to file
    out_path = "docs/presentation/visual_qa.md"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(report_lines))
        
    print(f"QA report written to {out_path}")

if __name__ == "__main__":
    pptx_path = "docs/presentation/YouCode_AI_Platform_Soutenance_Modern.pptx"
    check_presentation(pptx_path)
